#!/usr/bin/env python3
"""Generate a .pptx file from a JSON spec.

Usage:
    python generate_pptx.py <spec.json> <output.pptx>

Spec schema: see ../references/spec_reference.md
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    sys.stderr.write(
        "error: python-pptx is not installed. Run: pip install python-pptx\n"
    )
    sys.exit(1)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

COLOR_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TITLE = RGBColor(0x0F, 0x1E, 0x3D)
COLOR_BODY = RGBColor(0x1F, 0x2A, 0x3D)
COLOR_ACCENT = RGBColor(0x1E, 0x6F, 0xD9)
COLOR_MUTED = RGBColor(0x6B, 0x76, 0x85)
COLOR_ON_DARK = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_ON_DARK_MUTED = RGBColor(0xC8, 0xD3, 0xE3)

CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide, left, top, width, height, text,
    *, size=18, bold=False, color=COLOR_BODY,
    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _accent_bar(slide, left, top, width=Inches(0.18), height=Inches(0.6), color=COLOR_ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _slide_title_band(slide, text):
    _accent_bar(slide, Inches(0.6), Inches(0.55))
    _add_text(
        slide, Inches(0.95), Inches(0.45), Inches(11.5), Inches(0.85),
        text, size=32, bold=True, color=COLOR_TITLE,
    )


def _bullets(slide, items, left, top, width, height, *, size=20):
    if not items:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = COLOR_BODY


def _new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, COLOR_BG)
    return slide


# --- Handlers ------------------------------------------------------------

def _slide_title(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _add_text(
        slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5),
        spec.get("title", ""), size=54, bold=True, color=COLOR_TITLE,
    )
    if spec.get("subtitle"):
        _add_text(
            slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.8),
            spec["subtitle"], size=24, color=COLOR_MUTED,
        )
    _accent_bar(slide, Inches(0.8), Inches(4.95), width=Inches(2.0), height=Inches(0.1))
    return slide


def _slide_section(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _set_bg(slide, COLOR_TITLE)
    _add_text(
        slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.2),
        spec.get("title", ""),
        size=44, bold=True, color=COLOR_ON_DARK,
    )
    if spec.get("subtitle"):
        _add_text(
            slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.6),
            spec["subtitle"], size=20, color=COLOR_ON_DARK_MUTED,
        )
    _accent_bar(slide, Inches(0.8), Inches(5.0), width=Inches(2.0), height=Inches(0.1))
    return slide


def _slide_bullets(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _slide_title_band(slide, spec.get("title", ""))
    _bullets(
        slide, spec.get("bullets", []),
        Inches(0.95), Inches(1.7), Inches(11.5), Inches(5.3),
    )
    return slide


def _slide_two_column(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _slide_title_band(slide, spec.get("title", ""))
    if spec.get("left_title"):
        _add_text(
            slide, Inches(0.95), Inches(1.7), Inches(5.5), Inches(0.5),
            spec["left_title"], size=22, bold=True, color=COLOR_ACCENT,
        )
    if spec.get("right_title"):
        _add_text(
            slide, Inches(6.85), Inches(1.7), Inches(5.5), Inches(0.5),
            spec["right_title"], size=22, bold=True, color=COLOR_ACCENT,
        )
    top = Inches(2.3) if spec.get("left_title") or spec.get("right_title") else Inches(1.7)
    _bullets(slide, spec.get("left", []), Inches(0.95), top, Inches(5.5), Inches(4.7))
    _bullets(slide, spec.get("right", []), Inches(6.85), top, Inches(5.5), Inches(4.7))
    return slide


def _slide_image(prs, spec, base_dir):
    slide = _new_slide(prs)
    _slide_title_band(slide, spec.get("title", ""))
    img_path = Path(spec["image_path"])
    if not img_path.is_absolute():
        img_path = base_dir / img_path
    if not img_path.is_file():
        raise FileNotFoundError(f"Image not found: {img_path}")
    slide.shapes.add_picture(
        str(img_path),
        Inches(1.5), Inches(1.7),
        width=Inches(10.3),
    )
    if spec.get("caption"):
        _add_text(
            slide, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.6),
            spec["caption"], size=14, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
        )
    return slide


def _slide_table(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _slide_title_band(slide, spec.get("title", ""))
    headers = spec.get("headers", [])
    rows = spec.get("rows", [])
    if not headers or not rows:
        return slide
    n_rows = len(rows) + 1
    n_cols = len(headers)
    height = Inches(min(5.0, 0.5 * n_rows + 0.4))
    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.95), Inches(1.7), Inches(11.4), height,
    ).table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TITLE
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(16)
                run.font.bold = True
                run.font.color.rgb = COLOR_ON_DARK
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row[:n_cols]):
            cell = table.cell(r_i, c_i)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = COLOR_BODY
    return slide


def _slide_chart(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _slide_title_band(slide, spec.get("title", ""))
    chart_kind = CHART_TYPES.get(spec.get("chart_type", "column"))
    if chart_kind is None:
        raise ValueError(
            f"Unknown chart_type: {spec.get('chart_type')!r}. "
            f"Use one of: {', '.join(CHART_TYPES)}"
        )
    data = CategoryChartData()
    data.categories = spec.get("categories", [])
    for series in spec.get("series", []):
        data.add_series(series.get("name", "Series"), series.get("values", []))
    slide.shapes.add_chart(
        chart_kind,
        Inches(1.0), Inches(1.7), Inches(11.3), Inches(5.3),
        data,
    )
    return slide


def _slide_quote(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _add_text(
        slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
        f"“{spec.get('text', '')}”",
        size=36, color=COLOR_TITLE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    if spec.get("attribution"):
        _add_text(
            slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.6),
            f"— {spec['attribution']}",
            size=18, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
        )
    return slide


def _slide_closing(prs, spec, _base_dir):
    slide = _new_slide(prs)
    _add_text(
        slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.5),
        spec.get("title", "Thank you"),
        size=54, bold=True, color=COLOR_TITLE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    if spec.get("subtitle"):
        _add_text(
            slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8),
            spec["subtitle"], size=22, color=COLOR_MUTED, align=PP_ALIGN.CENTER,
        )
    return slide


HANDLERS = {
    "title": _slide_title,
    "section": _slide_section,
    "bullets": _slide_bullets,
    "two_column": _slide_two_column,
    "image": _slide_image,
    "table": _slide_table,
    "chart": _slide_chart,
    "quote": _slide_quote,
    "closing": _slide_closing,
}


def build(spec, output, base_dir):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    if spec.get("title"):
        prs.core_properties.title = spec["title"]
    if spec.get("author"):
        prs.core_properties.author = spec["author"]

    slides_spec = spec.get("slides", [])
    for i, slide_spec in enumerate(slides_spec):
        slide_type = slide_spec.get("type", "bullets")
        handler = HANDLERS.get(slide_type)
        if handler is None:
            raise ValueError(
                f"Unknown slide type at index {i}: {slide_type!r}. "
                f"Valid types: {', '.join(HANDLERS)}"
            )
        slide = handler(prs, slide_spec, base_dir)
        notes = slide_spec.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return len(slides_spec)


def main(argv):
    if len(argv) != 3:
        sys.stderr.write("usage: generate_pptx.py <spec.json> <output.pptx>\n")
        return 2
    spec_path = Path(argv[1]).resolve()
    out_path = Path(argv[2]).resolve()
    if out_path.suffix.lower() != ".pptx":
        sys.stderr.write("error: output path must end in .pptx\n")
        return 2
    if not spec_path.is_file():
        sys.stderr.write(f"error: spec not found: {spec_path}\n")
        return 2
    spec = json.loads(spec_path.read_text())
    n = build(spec, out_path, spec_path.parent)
    print(f"Wrote {n} slides → {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
